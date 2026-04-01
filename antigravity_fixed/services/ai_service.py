"""
AI Service
----------
Three capabilities:
  1. generate_pitch_deck  — Aggregates project + team data, calls GPT-4o,
                            returns structured JSON (10 slides), renders to .pptx
  2. generate_lean_canvas — Fills the 9-block Lean Canvas from project context
  3. startup_readiness    — Scores the project 0–100 across 5 axes

All LLM calls use structured JSON output mode to guarantee parseable responses.
"""

import json
import uuid
import io
from typing import Optional
from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from config import settings
import structlog

log = structlog.get_logger()

_client: Optional[AsyncOpenAI] = None


def _openai() -> AsyncOpenAI:
    global _client
    if _client is None:
        _client = AsyncOpenAI(api_key=settings.OPENAI_API_KEY)
    return _client


# ── Pitch Deck ─────────────────────────────────────────────────────────────────

PITCH_DECK_PROMPT = """
You are a top-tier startup pitch advisor with experience coaching Y Combinator and Sequoia-backed founders.
Generate a structured 10-slide investor pitch deck as JSON.

Project context:
{context}

Return ONLY valid JSON with this exact structure:
{{
  "slides": [
    {{
      "slide_number": 1,
      "title": "Problem",
      "headline": "One powerful sentence about the problem",
      "bullets": ["bullet 1", "bullet 2", "bullet 3"],
      "visual_note": "Suggested visual or chart description"
    }}
  ]
}}

Slides must be: Problem, Solution, Market Size, Product, Business Model,
Traction & Milestones, Go-To-Market, Team, Financials (Ask), Vision.
Be specific, concise, and investor-ready. No filler content.
"""


async def generate_pitch_deck(project_data: dict, team_data: dict, milestones: list) -> dict:
    """
    Returns dict with keys: slides (list), raw_json (str)
    Raises ValueError if GPT response is not parseable.
    """
    context = (
        f"Project: {project_data.get('title')}\n"
        f"Problem: {project_data.get('problem_statement', 'N/A')}\n"
        f"Target Market: {project_data.get('target_market', 'N/A')}\n"
        f"Industry: {project_data.get('industry_vertical', 'N/A')}\n"
        f"Stage: {project_data.get('stage')}\n"
        f"Completed Milestones: {[m['title'] for m in milestones if m.get('status') == 'completed']}\n"
        f"Team Domains: {[m['domain'] for m in team_data.get('members', [])]}\n"
    )

    response = await _openai().chat.completions.create(
        model=settings.LLM_MODEL,
        max_tokens=2000,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": "You are a startup pitch advisor. Always respond with valid JSON only."},
            {"role": "user", "content": PITCH_DECK_PROMPT.format(context=context)},
        ],
    )

    raw = response.choices[0].message.content
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError as e:
        log.error("pitch_deck_parse_error", error=str(e))
        raise ValueError("LLM returned invalid JSON for pitch deck")

    return {"slides": parsed.get("slides", []), "raw_json": raw}


def render_pptx(slides: list, project_title: str) -> bytes:
    """
    Renders a list of slide dicts to a .pptx binary blob.
    Returns bytes ready for upload to S3/R2.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x0D, 0x0D, 0x1A)
    ACCENT = RGBColor(0x6C, 0x63, 0xFF)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREY = RGBColor(0xCC, 0xCC, 0xCC)

    blank_layout = prs.slide_layouts[6]  # Completely blank

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Slide number pill
        _add_textbox(slide, f"{slide_data.get('slide_number', '')}", 0.2, 0.2, 0.6, 0.4,
                     font_size=10, color=GREY, bold=False)

        # Slide title tag (accent color)
        _add_textbox(slide, slide_data.get("title", "").upper(), 0.5, 0.3, 4, 0.4,
                     font_size=11, color=ACCENT, bold=True)

        # Headline
        _add_textbox(slide, slide_data.get("headline", ""), 0.5, 0.9, 12, 1.1,
                     font_size=28, color=WHITE, bold=True)

        # Bullets
        bullets = slide_data.get("bullets", [])
        for i, bullet in enumerate(bullets[:4]):
            _add_textbox(slide, f"→  {bullet}", 0.5, 2.2 + i * 0.9, 11, 0.8,
                         font_size=16, color=GREY, bold=False)

        # Visual note (bottom right, muted)
        visual = slide_data.get("visual_note", "")
        if visual:
            _add_textbox(slide, f"[{visual}]", 0.5, 6.7, 12, 0.5,
                         font_size=9, color=RGBColor(0x66, 0x66, 0x88), bold=False)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_textbox(slide, text, left, top, width, height, font_size, color, bold):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold


# ── Lean Canvas ────────────────────────────────────────────────────────────────

LEAN_CANVAS_PROMPT = """
You are an expert startup advisor. Fill in the Lean Canvas for this project.
Return ONLY valid JSON with exactly these 9 keys:
problem, solution, unique_value_proposition, unfair_advantage,
customer_segments, key_metrics, channels, cost_structure, revenue_streams.

Each value should be 2-3 concise bullet points as a list of strings.

Project context:
{context}
"""


async def generate_lean_canvas(project_data: dict) -> dict:
    context = (
        f"Title: {project_data.get('title')}\n"
        f"Problem: {project_data.get('problem_statement', 'N/A')}\n"
        f"Target Market: {project_data.get('target_market', 'N/A')}\n"
        f"Description: {project_data.get('description', '')[:500]}\n"
    )
    response = await _openai().chat.completions.create(
        model=settings.LLM_MODEL,
        max_tokens=1200,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": "You are a startup advisor. Respond with valid JSON only."},
            {"role": "user", "content": LEAN_CANVAS_PROMPT.format(context=context)},
        ],
    )
    raw = response.choices[0].message.content
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        raise ValueError("LLM returned invalid JSON for lean canvas")


# ── Startup Readiness Score ────────────────────────────────────────────────────

READINESS_PROMPT = """
Evaluate this student startup project across 5 axes.
Return ONLY valid JSON:
{{
  "problem_clarity": <0-20>,
  "market_size": <0-20>,
  "solution_viability": <0-20>,
  "team_completeness": <0-20>,
  "execution_evidence": <0-20>,
  "total": <0-100>,
  "summary": "2-3 sentence honest assessment",
  "top_risks": ["risk1", "risk2", "risk3"],
  "next_actions": ["action1", "action2", "action3"]
}}

Project data:
{context}
"""


async def startup_readiness_score(project_data: dict, team_data: dict, milestone_count: int) -> dict:
    context = (
        f"Title: {project_data.get('title')}\n"
        f"Problem: {project_data.get('problem_statement', 'N/A')}\n"
        f"Market: {project_data.get('target_market', 'N/A')}\n"
        f"Stage: {project_data.get('stage')}\n"
        f"Team size: {len(team_data.get('members', []))}\n"
        f"Team domains: {[m['domain'] for m in team_data.get('members', [])]}\n"
        f"Completed milestones: {milestone_count}\n"
    )
    response = await _openai().chat.completions.create(
        model=settings.LLM_MODEL,
        max_tokens=800,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": "You are a startup evaluator. Respond with valid JSON only."},
            {"role": "user", "content": READINESS_PROMPT.format(context=context)},
        ],
    )
    raw = response.choices[0].message.content
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        raise ValueError("LLM returned invalid JSON for readiness score")


# ── AI Co-Founder Chat ─────────────────────────────────────────────────────────

COFOUNDER_CHAT_SYSTEM = """\
You are an AI co-founder advisor embedded inside the Antigravity platform — a cross-disciplinary student startup incubator.
You have full context of this project's current state, team, and milestones.

Adapt your advice based on the caller's domain:
- engineering  → architecture decisions, tech stack choices, API design, scalability
- design       → UX flows, user research approach, visual direction, design system
- business     → GTM strategy, revenue model, market sizing, investor narrative

Be direct, specific, and actionable. Reference the actual project data in every answer.
Never give generic startup advice — always tie it to the real context below.

Project: {project_context}
Team: {team_context}
Milestones: {milestone_context}
Caller domain: {caller_domain}
"""


async def cofounder_chat(
    project_data: dict,
    team_data: dict,
    milestones: list,
    caller_domain: str,
    message: str,
    conversation_history: list,
) -> str:
    """
    Multi-turn project-scoped AI co-founder.
    Keeps last 10 turns of history to stay within context window.
    Raises ValueError if the LLM call fails.
    """
    project_context = (
        f"{project_data['title']} | Stage: {project_data['stage']} | "
        f"Problem: {project_data.get('problem_statement', 'N/A')} | "
        f"Market: {project_data.get('target_market', 'N/A')}"
    )
    team_context = (
        ", ".join(f"{m['name']} ({m['domain']})" for m in team_data.get("members", []))
        or "No team members yet"
    )
    milestone_context = (
        " | ".join(f"{m['title']} [{m['status']}]" for m in milestones)
        or "No milestones set"
    )

    system_msg = COFOUNDER_CHAT_SYSTEM.format(
        project_context=project_context,
        team_context=team_context,
        milestone_context=milestone_context,
        caller_domain=caller_domain,
    )

    messages = [{"role": "system", "content": system_msg}]
    # Include last 10 turns to avoid ballooning context
    for h in conversation_history[-10:]:
        messages.append({"role": h["role"], "content": h["content"]})
    messages.append({"role": "user", "content": message})

    response = await _openai().chat.completions.create(
        model=settings.LLM_MODEL,
        max_tokens=800,
        messages=messages,
    )
    return response.choices[0].message.content


# ── Smart Roadmap Generator ────────────────────────────────────────────────────

ROADMAP_PROMPT = """\
You are a senior product manager running a 6-week sprint for a student startup.

Rules:
1. Read each member's availability_hours — NEVER assign more than their weekly hours
2. Assign each deliverable to the right domain owner (engineering / design / business)
3. Mark critical-path items — milestone B cannot start until milestone A completes
4. If the team is missing a domain critical for this stage, set missing_role_alert
5. Each sprint covers 2 weeks (3 sprints total = 6 weeks)

Project: {project_context}
Team with availability: {team_with_hours}
Existing milestones: {milestones}
Stage: {stage}

Return ONLY valid JSON matching this exact structure:
{{
  "sprints": [
    {{
      "sprint_number": 1,
      "title": "Sprint title",
      "week_range": "Weeks 1-2",
      "goal": "Sprint goal in one sentence",
      "milestones": [
        {{
          "title": "Deliverable title",
          "owner_domain": "engineering",
          "estimated_hours": 8,
          "is_critical_path": true,
          "depends_on": []
        }}
      ]
    }}
  ],
  "critical_path_summary": "Which items are blocking and why",
  "bandwidth_warning": null,
  "missing_role_alert": null,
  "definition_of_done": "What success looks like at end of week 6"
}}
"""


async def generate_roadmap(project_data: dict, team_data: dict, milestones: list) -> dict:
    """
    Generates a 3-sprint roadmap that respects availability hours per member.
    Raises ValueError on LLM parse failure.
    """
    team_with_hours = (
        [
            f"{m['name']} ({m['domain']}, {m.get('availability_hours', 10)}h/week)"
            for m in team_data.get("members", [])
        ]
        or ["No team members yet"]
    )
    project_context = (
        f"{project_data['title']} | "
        f"Problem: {project_data.get('problem_statement', 'N/A')} | "
        f"Market: {project_data.get('target_market', 'N/A')}"
    )
    milestone_titles = [m["title"] for m in milestones] or ["No milestones set"]

    response = await _openai().chat.completions.create(
        model=settings.LLM_MODEL,
        max_tokens=2000,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": "You are a startup PM. Respond with valid JSON only."},
            {
                "role": "user",
                "content": ROADMAP_PROMPT.format(
                    project_context=project_context,
                    team_with_hours=team_with_hours,
                    milestones=milestone_titles,
                    stage=project_data["stage"],
                ),
            },
        ],
    )
    raw = response.choices[0].message.content
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        log.error("roadmap_parse_error", error=str(e))
        raise ValueError("LLM returned invalid JSON for roadmap")


# ── Idea Validator ─────────────────────────────────────────────────────────────

VALIDATE_PROMPT = """\
You are a startup accelerator judge evaluating a student team's idea.
Be honest, specific, and reference the actual idea details — no generic advice.

Idea:
Title: {title}
Problem: {problem}
Market: {market}
Industry: {industry}
Description: {description}

Return ONLY valid JSON with this exact structure:
{{
  "viability_score": <0-100>,
  "overall_grade": "A",
  "originality_score": <0-100>,
  "cross_disciplinary_need_score": <0-100>,
  "recommended_team_composition": {{
    "engineering_pct": <integer 0-100>,
    "design_pct": <integer 0-100>,
    "business_pct": <integer 0-100>
  }},
  "mvp_suggestion": "Concrete 1-sentence MVP that could be built in 4 weeks",
  "green_flags": ["specific strength 1", "specific strength 2", "specific strength 3"],
  "red_flags": ["specific risk 1", "specific risk 2"],
  "pivot_suggestions": ["pivot idea 1", "pivot idea 2"],
  "first_customer_hypothesis": "Who is the very first user and why they would pay",
  "verdict": "2-sentence honest overall assessment"
}}

overall_grade must be one of: A, B, C, D, F
engineering_pct + design_pct + business_pct must equal 100.
"""


async def validate_idea(
    title: str,
    problem: str,
    market: str,
    industry: str,
    description: str,
) -> dict:
    """
    Pre-flight idea validation — called before the project row is created.
    Returns a rich scorecard with grade, flags, team composition, and MVP suggestion.
    Raises ValueError on LLM parse failure.
    """
    response = await _openai().chat.completions.create(
        model=settings.LLM_MODEL,
        max_tokens=1000,
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": "You are a startup evaluator. Respond with valid JSON only."},
            {
                "role": "user",
                "content": VALIDATE_PROMPT.format(
                    title=title,
                    problem=problem,
                    market=market,
                    industry=industry or "Not specified",
                    description=description or "Not provided",
                ),
            },
        ],
    )
    raw = response.choices[0].message.content
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        log.error("idea_validate_parse_error", error=str(e))
        raise ValueError("LLM returned invalid JSON for idea validation")
